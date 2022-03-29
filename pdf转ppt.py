# -*- coding: utf-8 -*-
from pptx import Presentation
from pdf2image import convert_from_path, convert_from_bytes
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
from sys import argv
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path

fff=argv[1]
pp=Path.cwd()
fn=fff+'.pdf'
filename=pp/fn
print("Creating %s" % fff)
prs = Presentation()
width,height=argv[3].split('x')
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

pages = convert_from_path(filename,dpi=int(argv[2]), size=(int(width)*100,int(height)*100))
jpgs=pp/'jpgs'
if not jpgs.exists():
    jpgs.mkdir()
for index, page in enumerate(pages):
    name=fff+"-(%d).png" % index
    jpg_file =jpgs/name
#    print(jpg_file)
    page.save(jpg_file, 'PNG')


    image = Image.open(jpg_file)
    height = image.height
    width = image.width
#
    if height > width:
        adjusted = image.rotate(270, expand=True)
        adjusted.save(jpg_file)
#
#
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    left = top = 0
    jpg_file=str(jpg_file)
    print(jpg_file)
    slide.shapes.add_picture(jpg_file, left,top,height = prs.slide_height)



pptname='%s.pptx' % fff
prs.save(pp/pptname)



print("Saved")